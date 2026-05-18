"""
text_utils.py — Extract and chunk text from PDF, PPTX, DOCX
"""
from __future__ import annotations
import io
import re
import unicodedata
from pypdf import PdfReader

CHUNK_WORDS   = 200
OVERLAP_WORDS = 40


def clean_text(raw: str) -> str:
    text = unicodedata.normalize("NFKC", raw)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    pages  = []
    for page in reader.pages:
        t = page.extract_text()
        if t and t.strip():
            pages.append(t.strip())
    return clean_text(" ".join(pages))


def extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs   = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(r.text for r in para.runs).strip()
                    if line:
                        slide_text.append(line)
        if slide_text:
            parts.append(f"[Slide {i}] " + " ".join(slide_text))
    return clean_text("\n".join(parts))


def extract_docx(data: bytes) -> str:
    from docx import Document
    doc  = Document(io.BytesIO(data))
    text = []
    for p in doc.paragraphs:
        if p.text.strip():
            text.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text.append(cell.text.strip())
    return clean_text("\n".join(text))


def chunk_text(text: str, doc_id: str) -> list:
    words = text.split()
    if not words:
        return []

    chunks = []
    start  = 0
    idx    = 0

    while start < len(words):
        end   = start + CHUNK_WORDS
        chunk = " ".join(words[start:end]).strip()

        if len(words[start:end]) >= 15:  # skip tiny fragments
            chunks.append({
                "chunk_id":    f"doc{doc_id}_c{idx}",
                "text":        chunk,
                "chunk_index": idx,
            })
            idx += 1

        if end >= len(words):
            break
        start += CHUNK_WORDS - OVERLAP_WORDS

    return chunks